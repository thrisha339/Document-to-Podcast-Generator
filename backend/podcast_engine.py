import io
import os
from dotenv import load_dotenv
load_dotenv()
import re
import time
import base64
import struct
import logging
import textwrap
from concurrent.futures import ThreadPoolExecutor

import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from groq import Groq
from elevenlabs import ElevenLabs

# Logging 
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  [%(levelname)s]  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("PodcastGen")
try:
    import pytesseract
    from pdf2image import convert_from_bytes
    OCR_AVAILABLE = True
    log.info("OCR support: ENABLED")
except ImportError:
    OCR_AVAILABLE = False
    log.warning("OCR support: DISABLED — install tesseract + pdf2image for scanned PDFs")

# API clients (keys read from environment)
GROQ_API_KEY       = os.environ.get("GROQ_API_KEY", "")
ELEVENLABS_API_KEY = os.environ.get("ELEVENLABS_API_KEY", "")
if not GROQ_API_KEY:
    raise EnvironmentError("Missing GROQ_API_KEY — add it to your .env file.")
if not ELEVENLABS_API_KEY:
    raise EnvironmentError("Missing ELEVENLABS_API_KEY — add it to your .env file.")
groq_client       = Groq(api_key=GROQ_API_KEY)
elevenlabs_client = ElevenLabs(api_key=ELEVENLABS_API_KEY)

log.info("ElevenLabs key loaded: %s", ELEVENLABS_API_KEY[:8])
try:
    voices = elevenlabs_client.voices.get_all()
    log.info("ElevenLabs connected successfully")
except Exception as e:
    log.error("ElevenLabs error: %s", e)

GROQ_MODEL   = "llama-3.3-70b-versatile"
VOICE_ALEX   = "nPczCjzI2devNBz1zQrb"
VOICE_JORDAN = "cgSgspJ2msm6clMCkdW9"
SAMPLE_RATE  = 22050

# Text extraction
def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[1].lower()
    log.info("Extracting text from .%s (%.1f KB)...", ext, len(file_bytes) / 1024)
    if ext == "pdf":
        text = _extract_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        doc  = DocxDocument(io.BytesIO(file_bytes))
        text = "\n".join(p.text for p in doc.paragraphs)
    elif ext in ("ppt", "pptx"):
        prs  = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append(para.text)
        text = "\n".join(lines)
    elif ext == "txt":
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: .{ext}")
    cleaned = _clean_text(text)
    log.info("Extracted %d chars / ~%d words.", len(cleaned), len(cleaned.split()))
    if not cleaned:
        raise ValueError(
            "No text could be extracted from the file. "
            "If it is a scanned PDF, install tesseract and pdf2image."
        )
    return cleaned
def _extract_pdf(file_bytes: bytes) -> str:
    reader      = pypdf.PdfReader(io.BytesIO(file_bytes))
    total_pages = len(reader.pages)
    log.info("PDF has %d page(s).", total_pages)
    pages_text, empty_pages = [], 0
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if not page_text.strip():
            empty_pages += 1
        else:
            pages_text.append(page_text)
    if total_pages > 0 and empty_pages / total_pages <= 0.5:
        return "\n".join(pages_text)
    if not OCR_AVAILABLE:
        raise ValueError(
            "This PDF appears to be scanned (image-only). "
            "Install tesseract-ocr and pdf2image to enable OCR support."
        )
    return _ocr_pdf(file_bytes, total_pages)
def _ocr_pdf(file_bytes: bytes, total_pages: int) -> str:
    dpi = 150 if total_pages > 30 else 200
    log.info("OCR: %d pages at %d DPI...", total_pages, dpi)
    t0     = time.time()
    images = convert_from_bytes(file_bytes, dpi=dpi, fmt="jpeg", thread_count=4)
    results: dict[int, str] = {}
    def _ocr_page(args):
        idx, img = args
        try:
            return idx, pytesseract.image_to_string(img, lang="eng", config="--oem 3 --psm 3")
        except Exception:
            return idx, ""
    with ThreadPoolExecutor(max_workers=4) as ex:
        for idx, text in ex.map(_ocr_page, enumerate(images)):
            results[idx] = text
    combined = "\n".join(results.get(i, "") for i in range(len(images)))
    log.info("OCR done: %d chars in %.1fs.", len(combined), time.time() - t0)
    if not combined.strip():
        raise ValueError("OCR extracted no text. The PDF may be very low quality.")
    return combined
def _clean_text(raw: str) -> str:
    text = re.sub(r"[ \t]+", " ", raw)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[^\x20-\x7E\n]", "", text)
    return text.strip()

_BOILERPLATE_LINE_PATTERNS = [
    r"^page\s*(no|number)?[\s.:]*\d*$", r"^-\s*\d+\s*-$",
    r"^\d+\s*\|?\s*p\s*a\s*g\s*e$",
    r"signature\s*of\s*(the\s*)?(faculty|examiner|supervisor|guide|hod)",
    r"^(faculty|examiner|hod|supervisor|guide)\s*(signature)?[\s.:_-]*$",
    r"verified\s*by", r"approved\s*by",
    r"roll\s*(no|number)\s*[:\-.]?\s*\d*",
    r"register\s*(no|number)\s*[:\-.]?\s*\d*",
    r"^(table\s*of\s*contents?|index|contents?)$",
    r"(university|institute|college|department)\s*of\s*(technology|science|engineering|arts)",
    r"(autonomous|affiliated\s*to|accredited\s*by)",
    r"(naac|nba|ugc|aicte)\s*(accredited|approved|recognized)?",
    r"(academic\s*year|batch)\s*[:\-.]?\s*\d{4}",
    r"^about\s*(the\s*)?(author|writer)s?$",
    r"all\s*rights?\s*reserved", r"copyright\s*\©?\s*\d{4}",
    r"isbn\s*[:\-.]?\s*[\d\-]+",
    r"www\.[a-z0-9\-]+\.[a-z]{2,}", r"^https?://",
    r"^[_\-=*#~.]{3,}$", r"^[\W\s]{0,3}$",
]
_BOILERPLATE_LINE_RE = re.compile("|".join(_BOILERPLATE_LINE_PATTERNS), re.IGNORECASE)
_BOILERPLATE_PARA_RE = re.compile(
    r"(table\s*of\s*contents?|pin\s*code|phone\s*no|fax\s*no)", re.IGNORECASE
)
LLM_FILTER_PROMPT = """\
Is this text chunk MEANINGFUL CONTENT (concepts, data, arguments, findings) worth discussing in a podcast?
Or is it BOILERPLATE (table of contents, author bio, cover page, signatures, page numbers, copyright notices)?
Reply ONLY with: KEEP or DISCARD

CHUNK:
{chunk}
"""
def _regex_filter(text: str) -> str:
    paragraphs = re.split(r"\n{2,}", text)
    kept = []
    for para in paragraphs:
        if _BOILERPLATE_PARA_RE.search(para):
            continue
        lines = para.splitlines()
        good  = [l for l in lines if not _BOILERPLATE_LINE_RE.match(l.strip())]
        if len(good) < 2 and len(lines) > 3:
            continue
        cleaned = "\n".join(good).strip()
        if cleaned:
            kept.append(cleaned)
    return "\n\n".join(kept)
def _llm_filter_chunks(chunks: list[str]) -> list[str]:
    kept = []
    for chunk in chunks:
        if len(chunk.split()) < 20:
            continue
        try:
            v = _groq(LLM_FILTER_PROMPT.format(chunk=chunk[:1500]), max_tokens=5).strip().upper()
            if v.startswith("KEEP"):
                kept.append(chunk)
        except Exception:
            kept.append(chunk)  
    return kept

def filter_boilerplate(text: str) -> str:
    after_regex = _regex_filter(text) or text
    candidates  = [
        c.strip()
        for c in textwrap.wrap(after_regex, width=1500, break_long_words=False)
        if c.strip()
    ]
    kept = _llm_filter_chunks(candidates)
    return "\n\n".join(kept) if kept else after_regex

# Chunking + summarisation
def chunk_text(text: str, size: int = 2000) -> list[str]:
    chunks = [c.strip() for c in textwrap.wrap(text, width=size, break_long_words=False) if c.strip()]
    log.info("Chunked into %d pieces.", len(chunks))
    return chunks

def _groq(prompt: str, max_tokens: int = 1024, retries: int = 3) -> str:
    for attempt in range(1, retries + 1):
        try:
            r = groq_client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=max_tokens,
                temperature=0.7,
            )
            return r.choices[0].message.content.strip()
        except Exception as e:
            log.warning("Groq attempt %d/%d: %s", attempt, retries, e)
            if attempt < retries:
                time.sleep(2 * attempt)
    raise RuntimeError("Groq API failed after all retries.")

def summarise_chunks(chunks: list[str]) -> list[str]:
    summaries = []
    for i, chunk in enumerate(chunks, 1):
        s = _groq(f"Summarise in 100-150 words, preserving key ideas:\n\n{chunk}")
        summaries.append(s)
        if i % 5 == 0:
            log.info("  Summarised %d/%d chunks.", i, len(chunks))
    return summaries

def hierarchical_summarise(summaries: list[str]) -> list[str]:
    if len(summaries) <= 30:
        return summaries
    batches = [summaries[i:i + 7] for i in range(0, len(summaries), 7)]
    return [
        _groq("Combine into ONE paragraph of 80-100 words:\n\n" + "\n\n".join(f"- {s}" for s in b))
        for b in batches
    ]
# Script generation
SCRIPT_PROMPT = """\
Write a natural podcast conversation between Alex and Jordan using these summaries.

RULES:
- Every line: Alex: <text>   OR   Jordan: <text>
- No narration or stage directions
- Short, conversational sentences
- ~{target} words total
- Cover all key points

SUMMARIES:
{summaries}

BEGIN:
"""
LINE_RE = re.compile(r"^(Alex|Jordan)\s*:\s*(.+)$", re.IGNORECASE)
def generate_script(summaries: list[str]) -> str:
    combined = "\n\n".join(summaries)
    target   = max(300, min(sum(len(s) for s in summaries) // 20, 2400))
    if len(combined) > 40_000:
        combined = combined[:40_000] + "\n[truncated]"

    log.info("Generating script (~%d words)...", target)
    r = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role": "user", "content": SCRIPT_PROMPT.format(summaries=combined, target=target)}],
        max_tokens=min(int(target * 1.4) + 400, 8192),
        temperature=0.8,
    )
    return r.choices[0].message.content.strip()
def parse_script(script: str) -> list[dict]:
    dialogue, last = [], "Alex"
    for line in script.splitlines():
        line = line.strip()
        if not line:
            continue
        m       = LINE_RE.match(line)
        speaker = m.group(1).capitalize() if m else last
        text    = m.group(2).strip()     if m else line
        if not text:
            continue
        dialogue.append({
            "speaker":  speaker,
            "text":     text,
            "voice_id": VOICE_ALEX if speaker == "Alex" else VOICE_JORDAN,
        })
        last = speaker
    return dialogue
# Audio generation
def _split_tts(text: str, size: int = 250) -> list[str]:
    if len(text) <= size:
        return [text]
    sentences, chunks, cur = re.split(r"(?<=[.!?])\s+", text), [], ""
    for s in sentences:
        if len(cur) + len(s) + 1 <= size:
            cur = (cur + " " + s).strip()
        else:
            if cur:
                chunks.append(cur)
            cur = s[:size]
    if cur:
        chunks.append(cur)
    return chunks or [text[:size]]

def _tts(text: str, voice_id: str) -> bytes:
    return b"".join(
        elevenlabs_client.text_to_speech.convert(
            text=text,
            voice_id=voice_id,
            model_id="eleven_turbo_v2",
            output_format="pcm_22050",
        )
    )

def generate_audio(dialogue: list[dict]) -> bytes:
    segments = []
    pause    = b"\x00\x00" * int(SAMPLE_RATE * 0.4)
    for i, entry in enumerate(dialogue, 1):
        for sub in _split_tts(entry["text"]):
            try:
                segments.append(_tts(sub, entry["voice_id"]))
            except Exception as e:
                log.warning("TTS line %d: %s", i, e)
        segments.append(pause)

    pcm       = b"".join(segments)
    data_size = len(pcm)
    byte_rate = SAMPLE_RATE * 2
    header = struct.pack(
        "<4sI4s4sIHHIIHH4sI",
        b"RIFF", 36 + data_size, b"WAVE",
        b"fmt ", 16, 1, 1, SAMPLE_RATE, byte_rate, 2, 16,
        b"data", data_size,
    )
    duration = data_size / byte_rate
    log.info("Audio: %.1f seconds (%.1f min).", duration, duration / 60)
    return header + pcm

def run_pipeline(file_bytes: bytes, filename: str) -> dict:
    t0 = time.time()
    log.info("=== Pipeline START: %s ===", filename)
    text      = extract_text(file_bytes, filename)
    filtered  = filter_boilerplate(text)
    chunks    = chunk_text(filtered)
    if not chunks:
        raise ValueError("No meaningful content found after filtering.")

    summaries  = summarise_chunks(chunks)
    final_sums = hierarchical_summarise(summaries)
    script     = generate_script(final_sums)
    dialogue   = parse_script(script)

    if not dialogue:
        raise ValueError("Script parsing produced no dialogue.")

    wav       = generate_audio(dialogue)
    audio_b64 = base64.b64encode(wav).decode()

    log.info("=== Pipeline DONE in %.1fs ===", time.time() - t0)
    return {"script": script, "audio": audio_b64}